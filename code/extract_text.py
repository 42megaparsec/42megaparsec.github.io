from pptx import Presentation
from deep_translator import GoogleTranslator

pure_name = "addlk_230823_statwork.pptx"

filename = "data/" + pure_name
filename2 = "data/translated_" + pure_name

translator = GoogleTranslator(source='ko', target='en')

prs = Presentation(filename)

for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            original = shape.text.strip()
            # 이미 영어인 경우는 그대로 둔다
            if all(ord(c) < 128 for c in original):
                continue
            # 번역
            translated = translator.translate(original)
            shape.text = translated
            print (f"Original: {original}")
            print (f"Translated: {translated}")

prs.save(filename2)